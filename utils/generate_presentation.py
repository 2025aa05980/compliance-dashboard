from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

OUT = "IAM_Compliance_Dashboard_Presentation.pptx"

NAVY = RGBColor(22, 43, 67)
TEAL = RGBColor(0, 145, 150)
BLUE = RGBColor(42, 120, 214)
GREEN = RGBColor(12, 163, 12)
AMBER = RGBColor(250, 178, 25)
RED = RGBColor(208, 59, 59)
INK = RGBColor(33, 45, 58)
MUTED = RGBColor(93, 108, 123)
PALE = RGBColor(243, 247, 248)
WHITE = RGBColor(255, 255, 255)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_bg(slide, color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    slide.shapes._spTree.remove(shape._element)
    slide.shapes._spTree.insert(2, shape._element)


def textbox(slide, text, x, y, w, h, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = "Aptos"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box


def title(slide, heading, kicker=None):
    textbox(slide, heading, 0.65, 0.45, 12.0, 0.55, 27, NAVY, True)
    if kicker:
        textbox(slide, kicker.upper(), 0.68, 0.18, 12.0, 0.2, 8, TEAL, True)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.68), Inches(1.12), Inches(0.72), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()


def footer(slide, number):
    textbox(slide, "IAM Compliance Dashboard", 0.68, 7.12, 4.0, 0.18, 8, MUTED)
    textbox(slide, str(number), 12.2, 7.12, 0.45, 0.18, 8, MUTED, align=PP_ALIGN.RIGHT)


def card(slide, x, y, w, h, heading, body, accent=TEAL, body_size=13):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PALE
    shape.line.color.rgb = RGBColor(220, 229, 232)
    accent_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    accent_shape.fill.solid()
    accent_shape.fill.fore_color.rgb = accent
    accent_shape.line.fill.background()
    textbox(slide, heading, x + 0.25, y + 0.16, w - 0.45, 0.34, 15, NAVY, True)
    textbox(slide, body, x + 0.25, y + 0.58, w - 0.45, h - 0.72, body_size, INK)


def bullet_box(slide, items, x, y, w, h, size=15, color=INK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    for index, item in enumerate(items):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(10)
        p.bullet = True
    return box

# 1 Cover
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.5), 0, Inches(3.833), prs.slide_height)
shape.fill.solid(); shape.fill.fore_color.rgb = TEAL; shape.line.fill.background()
textbox(slide, "IAM & PAM", 0.85, 1.35, 6.8, 0.45, 15, RGBColor(133, 224, 220), True)
textbox(slide, "Compliance Dashboard", 0.8, 1.9, 8.6, 1.25, 40, WHITE, True)
textbox(slide, "A practical prototype for identity, privileged access,\nand asset compliance reporting", 0.85, 3.45, 7.1, 0.85, 20, RGBColor(218, 229, 236))
textbox(slide, "Technology overview | operating model | deployment | security roadmap", 0.85, 5.75, 7.2, 0.35, 12, RGBColor(183, 201, 213))
textbox(slide, "01", 10.15, 1.0, 2.0, 0.8, 30, WHITE, True)
textbox(slide, "DEMO READY", 10.18, 2.0, 2.0, 0.3, 10, NAVY, True)
textbox(slide, "Dash\nPython\nRender", 10.15, 3.0, 2.1, 1.7, 23, WHITE, True)

# 2 Why
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide); title(slide, "Why this dashboard exists", "The problem")
card(slide, 0.7, 1.65, 3.75, 3.7, "Fragmented visibility", "IAM data is often spread across spreadsheets, asset inventories, PAM tools, and application records.\n\nThis dashboard brings the posture into one operating view.", BLUE)
card(slide, 4.8, 1.65, 3.75, 3.7, "Risk-focused reporting", "Leaders need a fast answer to:\n\n• What is compliant?\n• What is exposed?\n• What needs remediation?\n• Which teams own the action?", AMBER)
card(slide, 8.9, 1.65, 3.75, 3.7, "Repeatable evidence", "Consistent tables, charts, filters, and CSV/Excel exports support recurring reviews, audit preparation, and stakeholder conversations.", TEAL)
textbox(slide, "Outcome: a single, role-aware view of IAM and PAM compliance posture.", 1.0, 6.05, 11.3, 0.42, 20, NAVY, True, PP_ALIGN.CENTER); footer(slide, 2)

# 3 Stack
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide); title(slide, "Technology stack", "How it is built")
stack = [
    ("Python", "Core language", BLUE), ("Plotly Dash", "Reactive web UI", TEAL),
    ("Dash Bootstrap", "Responsive layout", NAVY), ("Pandas + NumPy", "Data processing", GREEN),
    ("Plotly", "Interactive charts", AMBER), ("openpyxl", "Excel export", RED),
    ("Flask + Gunicorn", "Web and production server", BLUE), ("Render + GitHub", "Deployment and source control", TEAL),
]
for i, (head, body, accent) in enumerate(stack):
    x = 0.8 + (i % 4) * 3.1; y = 1.55 + (i // 4) * 2.2
    card(slide, x, y, 2.7, 1.55, head, body, accent, 12)
textbox(slide, "Design principle: Python-native delivery with a replaceable data layer.", 1.0, 6.2, 11.2, 0.35, 17, NAVY, True, PP_ALIGN.CENTER); footer(slide, 3)

# 4 Architecture
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide); title(slide, "Application architecture", "From data to decision")
steps = [
    ("01", "CSV / JSON", "11 IAM and asset datasets\nField visibility configuration"),
    ("02", "Data layer", "Pandas loader\nNormalization and summaries"),
    ("03", "Dash UI", "Role-based pages\nCharts, tables, filters"),
    ("04", "Evidence", "CSV / Excel exports\nOperational review"),
]
for i, (num, head, body) in enumerate(steps):
    x = 0.75 + i * 3.1
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(2.0), Inches(0.65), Inches(0.65))
    circle.fill.solid(); circle.fill.fore_color.rgb = TEAL; circle.line.fill.background()
    textbox(slide, num, x, 2.0, 0.65, 0.65, 13, WHITE, True, PP_ALIGN.CENTER)
    card(slide, x + 0.05, 2.9, 2.6, 1.65, head, body, TEAL, 12)
    if i < 3:
        line = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 2.68), Inches(2.18), Inches(0.42), Inches(0.28))
        line.fill.solid(); line.fill.fore_color.rgb = RGBColor(180, 198, 203); line.line.fill.background()
textbox(slide, "The loader can later be replaced with SQLAlchemy, CyberArk APIs, or CMDB REST calls.", 1.0, 5.7, 11.2, 0.5, 16, MUTED, False, PP_ALIGN.CENTER); footer(slide, 4)

# 5 Capabilities
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide); title(slide, "What the prototype demonstrates", "Core capabilities")
card(slide, 0.75, 1.55, 3.8, 4.65, "Identity coverage", "Human accounts\nPrivileged accounts\nService accounts\nBot accounts\nAI agents", BLUE, 15)
card(slide, 4.78, 1.55, 3.8, 4.65, "Asset coverage", "Windows servers\nLinux servers\nNetwork devices\nVirtual / ESXi assets\nApplications\nBreak-glass resources", TEAL, 15)
card(slide, 8.8, 1.55, 3.8, 4.65, "Operating features", "Executive summary\nCompliance and risk charts\nRole-based navigation\nField-level visibility\nDynamic query\nCSV and Excel export", GREEN, 15); footer(slide, 5)

# 6 Roles
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide); title(slide, "Role-based operating model", "Who uses it")
roles = [
    ("Executive", "Enterprise posture and key findings", RED),
    ("Leadership", "Cross-domain compliance oversight", AMBER),
    ("Operations", "Detailed records and remediation", BLUE),
    ("Admin", "Field visibility and user registry", GREEN),
]
for i, (head, body, accent) in enumerate(roles):
    y = 1.55 + i * 1.2
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.1), Inches(y), Inches(11.1), Inches(0.9))
    shape.fill.solid(); shape.fill.fore_color.rgb = PALE; shape.line.color.rgb = RGBColor(220, 229, 232)
    marker = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.1), Inches(y), Inches(0.16), Inches(0.9))
    marker.fill.solid(); marker.fill.fore_color.rgb = accent; marker.line.fill.background()
    textbox(slide, head, 1.55, y + 0.12, 2.4, 0.3, 17, NAVY, True)
    textbox(slide, body, 4.0, y + 0.12, 7.6, 0.3, 15, INK)
textbox(slide, "Visibility is tailored by role; production authorization should be enforced server-side as well.", 1.1, 6.25, 11.0, 0.35, 15, MUTED, False, PP_ALIGN.CENTER); footer(slide, 6)

# 7 Metrics
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide); title(slide, "Executive view and metrics", "Decision support")
metrics = [("1,090", "simulated records", BLUE), ("11", "data domains", TEAL), ("13", "dashboard pages", GREEN), ("4", "user roles", AMBER)]
for i, (value, label, accent) in enumerate(metrics):
    x = 0.85 + i * 3.1
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.65), Inches(2.65), Inches(1.65))
    shape.fill.solid(); shape.fill.fore_color.rgb = PALE; shape.line.color.rgb = RGBColor(220, 229, 232)
    textbox(slide, value, x + 0.15, 1.88, 2.35, 0.55, 30, accent, True, PP_ALIGN.CENTER)
    textbox(slide, label, x + 0.15, 2.55, 2.35, 0.28, 12, NAVY, True, PP_ALIGN.CENTER)
card(slide, 1.05, 4.0, 5.25, 1.55, "Compliance posture", "Overall compliance, non-compliant records, exceptions, and not-assessed populations.", GREEN)
card(slide, 7.0, 4.0, 5.25, 1.55, "Control posture", "PAM onboarding, password management, authentication integration, and risk ratings.", BLUE)
footer(slide, 7)

# 8 Deployment
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide); title(slide, "Deployment model", "How to share the demo")
card(slide, 0.8, 1.55, 3.65, 3.75, "Source", "GitHub repository\n\n2025aa05980/compliance-dashboard\n\nVersioned source, requirements, data, and Render configuration.", NAVY, 14)
card(slide, 4.85, 1.55, 3.65, 3.75, "Build", "Render installs requirements.txt and runs:\n\nGunicorn app:server\n\nThe Dash WSGI server is exposed over HTTPS.", TEAL, 14)
card(slide, 8.9, 1.55, 3.65, 3.75, "Access", "Share the generated onrender.com URL with your office.\n\nThe free service may sleep when idle, so the first request can be slower.", BLUE, 14)
textbox(slide, "Recommended for prototype demonstrations; use private enterprise hosting for sensitive IAM data.", 1.0, 6.0, 11.3, 0.42, 17, RED, True, PP_ALIGN.CENTER); footer(slide, 8)

# 9 Security
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide); title(slide, "Security posture", "Important boundary")
card(slide, 0.75, 1.55, 5.7, 4.45, "Current prototype", "• Demo credentials in project data\n• Session-based browser state\n• CSV-backed sample data\n• UI role filtering\n• Suitable for non-sensitive demonstrations", AMBER, 15)
card(slide, 6.85, 1.55, 5.7, 4.45, "Production target", "• Entra ID / Okta SSO and MFA\n• Server-side signed sessions\n• Backend authorization on callbacks\n• Hashed credentials or no local passwords\n• Secure database / APIs\n• Audit logs, secrets management, private hosting", GREEN, 15)
textbox(slide, "Conclusion: Dash is a reasonable internal dashboard framework; this implementation requires hardening before restricted production use.", 1.0, 6.3, 11.3, 0.4, 15, NAVY, True, PP_ALIGN.CENTER); footer(slide, 9)

# 10 Roadmap
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, NAVY); title(slide, "Recommended next steps", "From prototype to production")
roadmap = [
    ("01", "Secure identity", "Integrate corporate SSO, MFA, and conditional access."),
    ("02", "Secure data", "Move from CSV files to approved database or governed APIs."),
    ("03", "Enforce controls", "Validate role permissions on every server-side action."),
    ("04", "Operate and audit", "Add logs, monitoring, backups, and review workflows."),
]
for i, (num, head, body) in enumerate(roadmap):
    y = 1.4 + i * 1.15
    textbox(slide, num, 1.0, y, 0.65, 0.48, 20, RGBColor(133, 224, 220), True)
    textbox(slide, head, 2.0, y, 2.6, 0.35, 18, WHITE, True)
    textbox(slide, body, 5.0, y, 6.8, 0.4, 15, RGBColor(218, 229, 236))
textbox(slide, "Thank you", 0.95, 6.25, 11.5, 0.42, 23, WHITE, True, PP_ALIGN.CENTER)

prs.save(OUT)
print(f"Created {OUT}")
